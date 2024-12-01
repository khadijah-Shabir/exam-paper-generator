import os
import streamlit as st
from groq import Groq
import pdfplumber
from pptx import Presentation
import pandas as pd
from docx import Document
import time
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas

# Streamlit configuration
st.set_page_config(
    page_title="Exam Paper Generator",
    page_icon="📝",
    layout="wide"
)

def generate_styled_pdf(title: str, content: str) -> BytesIO:
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)

    # Create custom styles
    styles = getSampleStyleSheet()
    
    # Title Style
    title_style = ParagraphStyle(
        'TitleStyle',
        parent=styles['Heading1'],
        fontSize=22,
        textColor=colors.HexColor("#2C3E50"),
        alignment=1,  # Center alignment
        spaceAfter=12,
        fontName='Helvetica-Bold'
    )
    
    # Question Header Style
    question_header_style = ParagraphStyle(
        'QuestionHeaderStyle',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=colors.HexColor("#34495E"),
        spaceBefore=12,
        spaceAfter=6,
        fontName='Helvetica-Bold'
    )
    
    # Normal Question Style
    question_style = ParagraphStyle(
        'QuestionStyle',
        parent=styles['BodyText'],
        fontSize=12,
        textColor=colors.black,
        spaceBefore=6,
        spaceAfter=6,
        fontName='Helvetica'
    )
    
    # Option Style (for MCQs)
    option_style = ParagraphStyle(
        'OptionStyle',
        parent=styles['BodyText'],
        fontSize=11,
        textColor=colors.HexColor("#2C3E50"),
        leftIndent=20,
        spaceBefore=3,
        spaceAfter=3,
        fontName='Helvetica'
    )
    
    # Prepare document elements
    elements = []

    # Add title
    elements.append(Paragraph(title, title_style))
    elements.append(Spacer(1, 0.25 * inch))

    # Process content with improved formatting
    lines = content.split('\n')
    for line in lines:
        line = line.strip()
        
        # Identify and style different types of content
        if line.startswith('###'):  # Section headers
            header = line.replace('###', '').strip()
            elements.append(Paragraph(header, question_header_style))
        elif line.startswith('Q'):  # Questions
            elements.append(Paragraph(line, question_style))
        elif line.startswith(('A)', 'B)', 'C)', 'D)')):  # MCQ Options
            elements.append(Paragraph(line, option_style))
        elif line:  # Other content
            elements.append(Paragraph(line, styles['BodyText']))
        
        # Add small spacing between elements
        elements.append(Spacer(1, 0.1 * inch))

    # Build PDF
    doc.build(elements)
    buffer.seek(0)
    return buffer

# Document processor class
class DocumentProcessor:
    def __init__(self, api_key: str):
        self.client = Groq(api_key=api_key)

    def extract_text(self, file) -> str:
        """Extract text from different document types."""
        try:
            if file.type == "application/pdf":
                with pdfplumber.open(file) as pdf:
                    return "\n".join(page.extract_text() or "" for page in pdf.pages)
            elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                doc = Document(file)
                return "\n".join(paragraph.text for paragraph in doc.paragraphs)
            elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                ppt = Presentation(file)
                return "\n".join(
                    shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")
                )
            elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                df = pd.read_excel(file)
                return df.to_string(index=False)
            elif file.type == "text/plain":
                return file.getvalue().decode("utf-8")
            else:
                return "Unsupported file format"
        except Exception as e:
            return f"Error processing file: {str(e)}"

    def generate_questions(self, content: str, question_type: str, num_questions: int, difficulty: str, specific_topic: str = None) -> str:
        """Generate questions based on content with type-specific prompts."""
        if not content or content.strip() == "":
            return "No content provided for question generation."

        # Separate prompts for each question type
        if question_type == "multiple choice questions":
            prompt = f"""
            MULTIPLE CHOICE QUESTIONS (MCQs) GENERATION INSTRUCTIONS:

            Context:
            - Source Content: {content[:4000]}
            {'- Focus Topic: ' + specific_topic if specific_topic else ''}
            - Difficulty Level: {difficulty}

            STRICT MCQ FORMATTING REQUIREMENTS:
            - Generate {num_questions} UNIQUE Multiple Choice Questions
            - Each question MUST have EXACTLY 4 options
            
            MCQ FORMAT:
            Q[number]. [Precise, knowledge-testing question]
            
            Options (EXACTLY 4, precisely formatted and new line after each option):
            A) [Option 1 then start a new line]
            B) [Option 2 then start a new line]
            C) [Option 3 then start a new line]
            D) [Option 4 then start a new line]
            
       

            CRITICAL GUIDELINES:
            - Derive questions ONLY from provided content
            - Ensure NO overlap between questions
            - Options must be academically credible
            - CORRECT answer must be unambiguously right
            - Maintain academic language
            - Complexity matches specified difficulty level
            """

        elif question_type == "short questions":
            prompt = f"""
            SHORT ANSWER QUESTIONS GENERATION INSTRUCTIONS:

            Context:
            - Source Content: {content[:4000]}
            {'- Focus Topic: ' + specific_topic if specific_topic else ''}
            - Difficulty Level: {difficulty}

            STRICT SHORT QUESTION FORMATTING REQUIREMENTS:
            - Generate {num_questions} UNIQUE Short Answer Questions
            - Each question requires a focused, concise response (2-3 sentences)
            
            SHORT QUESTION FORMAT:
            Q[number]. [Precise, concept-testing question requiring brief, specific answer]

            CRITICAL GUIDELINES:
            - Questions must be answerable using ONLY the provided content
            - Focus on key concepts, definitions, explanations
            - Avoid yes/no questions
            - Ensure questions test understanding, not mere recall
            - Each question should require analysis or explanation
            - Maintain academic rigor
            - Complexity matches specified difficulty level
            """

        elif question_type == "long questions":
            prompt = f"""
            LONG ANSWER QUESTIONS GENERATION INSTRUCTIONS:

            Context:
            - Source Content: {content[:4000]}
            {'- Focus Topic: ' + specific_topic if specific_topic else ''}
            - Difficulty Level: {difficulty}

            STRICT LONG QUESTION FORMATTING REQUIREMENTS:
            - Generate {num_questions} UNIQUE Comprehensive Questions
            - Each question requires an in-depth, multi-part response
            
            LONG QUESTION FORMAT:
            Q[number]. [Complex, analytical question requiring comprehensive explanation]


            CRITICAL GUIDELINES:
            - Questions must demand critical thinking
            - Require synthesis of information from content
            - Encourage analytical and evaluative responses
            - Include potential for original insight
            - Ensure questions are NOT simply information regurgitation
            - Complexity significantly higher than short questions
            - Match specified difficulty level precisely
            """

        try:
            response = self.client.chat.completions.create(
                messages=[
                    {"role": "system", "content": "You are a precise, academic exam question generator. Dont give the raw heading and data and give point to point data only in a good format and alignment"},
                    {"role": "user", "content": prompt}
                ],
                model="gemma2-9b-it",
                temperature=0.7  # Slight randomness to prevent repetition
            )
            return response.choices[0].message.content
        except Exception as e:
            return f"Error generating questions: {str(e)}"

# Main application remains the same as in previous version
def main():
    st.title("📝 Exam Paper Generator")
    st.write("Upload your documents to generate an exam paper with customized questions.")

    # API Key Management
    api_key = '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'

    # File uploader
    uploaded_files = st.file_uploader(
        "📎 Upload Documents",
        type=["pdf", "pptx", "docx", "xlsx", "txt"],
        accept_multiple_files=True
    )

    # Topic input
    specific_topic = st.text_input("🎯 Specify a Specific Topic (Optional)", 
                                   help="Enter a topic to focus the exam questions. Leave blank to use entire document content.")

    # Proceed only if files are uploaded
    if uploaded_files:
        try:
            # Initialize processor
            processor = DocumentProcessor(api_key)

            # Extract content from uploaded files
            combined_text = "\n\n".join(processor.extract_text(file) for file in uploaded_files)
            
            if not combined_text or combined_text.strip() == "":
                st.error("No text could be extracted from the uploaded documents.")
                return

            st.success("Documents uploaded and processed successfully!")

            # Question type selection
            st.write("### Select Question Types")
            mcq_selected = st.checkbox("Multiple Choice Questions (MCQs)")
            short_selected = st.checkbox("Short Questions")
            long_selected = st.checkbox("Long Questions")

            # Customize question generation
            question_settings = {}
            if mcq_selected:
                st.subheader("Settings for MCQs")
                num_mcqs = st.slider("Number of MCQs:", 1, 20, 5, key="num_mcqs")
                difficulty_mcqs = st.selectbox("Difficulty Level for MCQs:", ["Easy", "Medium", "Hard", "Mixed"], key="diff_mcqs")
                question_settings["mcq"] = {"num_questions": num_mcqs, "difficulty": difficulty_mcqs}

            if short_selected:
                st.subheader("Settings for Short Questions")
                num_short = st.slider("Number of Short Questions:", 1, 20, 5, key="num_short")
                difficulty_short = st.selectbox("Difficulty Level for Short Questions:", ["Easy", "Medium", "Hard", "Mixed"], key="diff_short")
                question_settings["short"] = {"num_questions": num_short, "difficulty": difficulty_short}

            if long_selected:
                st.subheader("Settings for Long Questions")
                num_long = st.slider("Number of Long Questions:", 1, 20, 5, key="num_long")
                difficulty_long = st.selectbox("Difficulty Level for Long Questions:", ["Easy", "Medium", "Hard", "Mixed"], key="diff_long")
                question_settings["long"] = {"num_questions": num_long, "difficulty": difficulty_long}

            # Generate and display questions
            if st.button("🚀 Generate Exam Paper"):
                if not question_settings:
                    st.warning("Please select at least one question type.")
                    return

                with st.spinner("Generating exam paper..."):
                    paper_content = []

                    # Generate questions for each selected type
                    for q_type, settings in question_settings.items():
                        question_type_name = {"mcq": "Multiple Choice Questions", "short": "Short Questions", "long": "Long Questions"}
                        questions = processor.generate_questions(
                            content=combined_text,
                            question_type=question_type_name[q_type].lower(),
                            num_questions=settings["num_questions"],
                            difficulty=settings["difficulty"],
                            specific_topic=specific_topic
                        )
                        paper_content.append(f"### {question_type_name[q_type]}\n\n{questions}")

                    # Combine all generated questions
                    combined_paper = "\n\n".join(paper_content)

                    # Display the generated questions
                    st.write("### Generated Exam Paper")
                    st.write(combined_paper)

                    # Provide download option
                    pdf_buffer = generate_styled_pdf("Exam Paper", combined_paper)
                    st.download_button(
                        label="📥 Download Paper as PDF",
                        data=pdf_buffer,
                        file_name="exam_paper.pdf",
                        mime="application/pdf"
                    )
        except Exception as e:
            st.error(f"An error occurred: {str(e)}")

if __name__ == "__main__":
    main()
